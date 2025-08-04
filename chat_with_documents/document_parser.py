import io
import os
from PIL import Image
from pptx import Presentation
from docx import Document
from transformers import BlipProcessor, BlipForConditionalGeneration
from typing import List
from pdf2image import convert_from_path
import fitz  # Only used for PDF image extraction

# Load BLIP-1 model and processor once
blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

def extract_images_from_pdf(pdf_path: str) -> List[Image.Image]:
    images = []
    doc = fitz.open(pdf_path)
    for page_index in range(len(doc)):
        for img_index, img in enumerate(doc.get_page_images(page_index)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            images.append(image)
    return images

def extract_images_from_pptx(pptx_path: str) -> List[Image.Image]:
    prs = Presentation(pptx_path)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                image_stream = shape.image.blob
                image = Image.open(io.BytesIO(image_stream)).convert("RGB")
                images.append(image)
    return images

def extract_images_from_docx(docx_path: str) -> List[Image.Image]:
    document = Document(docx_path)
    images = []
    for rel in document.part._rels:
        rel = document.part._rels[rel]
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            image = Image.open(io.BytesIO(image_data)).convert("RGB")
            images.append(image)
    return images

def generate_image_captions(images):
    captions = []
    for img in images:
        try:
            inputs = blip_processor(images=img, return_tensors="pt")
            out = blip_model.generate(**inputs)
            caption = blip_processor.decode(out[0], skip_special_tokens=True)
            captions.append(caption)
        except Exception as e:
            captions.append("⚠️ Captioning failed.")
    return captions

def process_images_from_file(file_path, extension):
    images = []
    if extension == "pdf":
        images = convert_from_path(file_path, dpi=300)
    elif extension == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    images.append(img)
    captions = generate_image_captions(images)
    return images, captions