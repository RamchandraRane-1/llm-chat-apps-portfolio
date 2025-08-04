# document_parser.py

import io
import time
from PIL import Image
from pptx import Presentation
from docx import Document
from pdf2image import convert_from_path
from typing import List, Tuple

def query_gemini(image: Image.Image, gemini_api_key: str, prompt: str = "Describe this image in detail.") -> str:
    import google.generativeai as genai

    if not gemini_api_key:
        return "⚠️ Gemini API key missing. Please provide a valid key."

    try:
        genai.configure(api_key=gemini_api_key)
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content([prompt, image])
        return response.text
    except Exception as e:
        return f"⚠️ Gemini API error: {e}"

def extract_images_from_pdf(pdf_path: str) -> List[Image.Image]:
    return convert_from_path(pdf_path, dpi=300)

def extract_images_from_pptx(pptx_path: str) -> List[Image.Image]:
    images = []
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    image_stream = shape.image.blob
                    image = Image.open(io.BytesIO(image_stream)).convert("RGB")
                    images.append(image)
                except Exception as e:
                    print(f"Error extracting image from PPTX: {e}")
    return images

def extract_images_from_docx(docx_path: str) -> List[Image.Image]:
    images = []
    document = Document(docx_path)
    for rel in document.part._rels:
        rel = document.part._rels[rel]
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                image = Image.open(io.BytesIO(image_data)).convert("RGB")
                images.append(image)
            except Exception as e:
                print(f"Error extracting image from DOCX: {e}")
    return images

def extract_images_from_file(file_path: str, extension: str) -> List[Image.Image]:
    if extension == ".pdf":
        return extract_images_from_pdf(file_path)
    elif extension == ".pptx":
        return extract_images_from_pptx(file_path)
    elif extension == ".docx":
        return extract_images_from_docx(file_path)
    return []

def process_images_from_file(file_path: str, extension: str, gemini_api_key: str) -> Tuple[List[Image.Image], List[str]]:
    images = extract_images_from_file(file_path, extension)
    captions = []

    for idx, img in enumerate(images):
        try:
            caption = query_gemini(img, gemini_api_key)
        except Exception as e:
            time.sleep(8)
            try:
                caption = query_gemini(img, gemini_api_key)
            except Exception as e:
                caption = f"⚠️ Gemini error: {e}"
        captions.append(caption)

    return images, captions
