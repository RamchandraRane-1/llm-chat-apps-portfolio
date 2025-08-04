# app.py

import os
import streamlit as st
from langchain.vectorstores import FAISS
from langchain.embeddings import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.chat_models import ChatOpenAI
from langchain.chains.question_answering import load_qa_chain
from langchain.docstore.document import Document
from document_parser import process_images_from_file
from PyPDF2 import PdfReader
from docx import Document as DocxReader
from pptx import Presentation
import tempfile

# ---- Sidebar Configuration ----
st.set_page_config(page_title="Gemini + ChatGPT File QA Bot")
st.set_page_config(page_title="Chatbot", layout="centered")

st.markdown("""
<style>
/* Gradient title with light/dark mode support */
.gradient-header {
    text-align: center;
    font-size: 40px;
    font-weight: bold;
    background: linear-gradient(45deg, #6A11CB, #2575FC);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 0px;
    animation: fadeIn 1.5s ease-in-out;
}

.subtext {
    text-align: center;
    font-size: 18px;
    color: var(--text-color-secondary, #888);
    margin-bottom: 30px;
    animation: fadeIn 2s ease-in-out;
}

.sidebar-title {
    font-weight: 700;
    font-size: 16px;
    color: white;
    margin-top: 20px;
    margin-bottom: 5px;
}

hr {
    border: none;
    border-top: 1px solid #ddd;
    margin: 10px 0;
}

@keyframes fadeIn {
    from { opacity: 0; transform: translateY(-8px); }
    to { opacity: 1; transform: translateY(0); }
}
</style>

<h1 class='gradient-header'>Chatbot</h1>
<p class='subtext'>Your AI assistant for documents & images</p>
""", unsafe_allow_html=True)

st.sidebar.markdown("<div class='sidebar-title'>API Keys</div>", unsafe_allow_html=True)
openai_api_key = st.sidebar.text_input("OpenAI API Key", type="password")
gemini_api_key = st.sidebar.text_input("Gemini API Key", type="password")
st.sidebar.markdown("<hr>", unsafe_allow_html=True)

if openai_api_key:
    st.session_state["API_Key"] = openai_api_key
if gemini_api_key:
    st.session_state["Gemini_API_Key"] = gemini_api_key

# ---- File Upload ----
st.sidebar.markdown("<div class='sidebar-title'>Upload File</div>", unsafe_allow_html=True)
uploaded_file = st.sidebar.file_uploader("Upload your document", type=["pdf", "docx", "pptx"])
st.markdown("""
<div style='text-align: center; font-size: 35px; font-weight: bold; margin-top: 10px; margin-bottom: 3px;'>
    What can I help with?
</div>
""", unsafe_allow_html=True)

query = st.text_input("", placeholder="e.g., Explain image 5 or What is the summary?", key="query_input")


if uploaded_file and openai_api_key and gemini_api_key:
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
        tmp_file.write(uploaded_file.read())
        tmp_file_path = tmp_file.name

    # ---- Extract Text ----
    raw_text = ""
    if file_extension == ".pdf":
        reader = PdfReader(tmp_file_path)
        for page in reader.pages:
            raw_text += page.extract_text() or ""
    elif file_extension == ".docx":
        doc = DocxReader(tmp_file_path)
        for para in doc.paragraphs:
            raw_text += para.text + "\n"
    elif file_extension == ".pptx":
        prs = Presentation(tmp_file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw_text += shape.text + "\n"

    # ---- Extract Image Captions via Gemini ----
    st.info("🔍 Analyzing images using Gemini...")
    images, captions = process_images_from_file(tmp_file_path, file_extension, gemini_api_key)

    if captions:
        raw_text += "\n\n📸 Image Descriptions from Gemini:\n"
        for idx, cap in enumerate(captions, 1):
            raw_text += f"\n🖼️ Image {idx}: {cap.strip()}"

    # ---- Show combined content (optional) ----
    with st.expander("📚 View Extracted Text + Image Descriptions"):
        st.text_area("Combined Context", raw_text, height=300)

    # ---- Process with LangChain ----
    text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=100)
    texts = text_splitter.split_text(raw_text)

    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
    db = FAISS.from_documents([Document(page_content=t) for t in texts], embeddings)

    # ---- Ask Question ----
    if query:
        docs = db.similarity_search(query)
        llm = ChatOpenAI(temperature=0, openai_api_key=openai_api_key)
        chain = load_qa_chain(llm, chain_type="stuff")
        response = chain.run(input_documents=docs, question=query)
        st.success(response)
